from fastapi import FastAPI, File, UploadFile, Form
from fastapi.responses import JSONResponse
import io
import gc
import os
from contextlib import asynccontextmanager
from PIL import Image
import pytesseract
import pypdfium2 as pdfium
import pandas as pd
from pptx import Presentation

models = {}

@asynccontextmanager
async def lifespan(app: FastAPI):
    print("--- STARTUP: LOADING MODELS ---")
    try:
        from surya.recognition import RecognitionPredictor
        from surya.detection import DetectionPredictor
        from surya.layout import LayoutPredictor
        from surya.foundation import FoundationPredictor
        
        foundation_predictor = FoundationPredictor()
        models['detection'] = DetectionPredictor()
        models['layout'] = LayoutPredictor(foundation_predictor)
        models['recognition'] = RecognitionPredictor(foundation_predictor)
        
        print("--- MODELS LOADED SUCCESSFULLY ---")
    except Exception as e:
        print(f"Error loading models: {e}")
    yield
    print("--- SHUTDOWN ---")

app = FastAPI(lifespan=lifespan)

@app.post("/ocr", response_class=JSONResponse)
async def process_document(
    file: UploadFile = File(...),
    engine: str = Form("surya") 
):
    content = await file.read()
    filename = file.filename.lower()
    response_data = {"filename": filename, "engine": engine, "pages": []}

    # --- HANDLER 1: EXCEL (.xlsx, .xls) ---
    if filename.endswith(('.xlsx', '.xls')):
        try:
            # Load into Pandas
            excel_file = io.BytesIO(content)
            xls = pd.ExcelFile(excel_file)
            
            for i, sheet_name in enumerate(xls.sheet_names):
                page_data = {"page": i+1, "sheet_name": sheet_name, "elements": []}
                
                # Read sheet
                df = pd.read_excel(xls, sheet_name=sheet_name)
                # Convert to string to avoid JSON errors with NaN/Dates
                text_content = df.to_string(index=False)
                
                # Add as Table
                page_data["elements"].append({
                    "type": "Table",
                    "content": text_content,
                    "raw_data": df.fillna("").to_dict(orient='records') # Bonus: Actual Data
                })
                response_data["pages"].append(page_data)
                
            return JSONResponse(content=response_data)
        except Exception as e:
             return JSONResponse(content={"error": f"Excel Error: {str(e)}"}, status_code=500)

    # --- HANDLER 2: POWERPOINT (.pptx) ---
    elif filename.endswith('.pptx'):
        try:
            ppt_file = io.BytesIO(content)
            prs = Presentation(ppt_file)
            
            for i, slide in enumerate(prs.slides):
                page_data = {"page": i+1, "elements": []}
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Try to guess type based on PPT placeholder names
                        elem_type = "Text"
                        if "title" in shape.name.lower():
                            elem_type = "Title"
                        elif "footer" in shape.name.lower():
                            elem_type = "PageFooter"
                        
                        page_data["elements"].append({
                            "type": elem_type,
                            "content": shape.text.strip()
                        })
                        
                response_data["pages"].append(page_data)
            
            return JSONResponse(content=response_data)
        except Exception as e:
            return JSONResponse(content={"error": f"PPTX Error: {str(e)}"}, status_code=500)

    # --- HANDLER 3: PDF / IMAGES (OCR) ---
    # (This is the existing Surya/Vision Logic)
    try:
        pil_images = []
        if filename.endswith(".pdf"):
            pdf = pdfium.PdfDocument(content)
            for i in range(len(pdf)):
                # High Speed Scale 1.5
                pil_images.append(pdf[i].render(scale=1.5).to_pil().convert("RGB"))
        else:
            pil_images.append(Image.open(io.BytesIO(content)).convert("RGB"))
            
        # ... [Paste the exact logic from the previous step here] ...
        # (I will include the condensed loop below for completeness)
        
        if engine == "surya":
            layout_predictor = models['layout']
            rec_predictor = models['recognition']
            det_predictor = models['detection']

            for i, image in enumerate(pil_images):
                print(f"--- Processing Page {i+1} ---")
                page_data = {"page": i+1, "elements": []}
                try:
                    layout_pred = layout_predictor([image])[0]
                    bboxes = sorted(layout_pred.bboxes, key=lambda x: x.bbox[1])
                    if not bboxes: raise Exception("No layout")

                    crops_to_process, crop_indices, temp_elements = [], [], []

                    for idx, box in enumerate(bboxes):
                        element = {"type": box.label, "content": ""}
                        if box.label in ["Picture", "Figure", "Table", "Image"]:
                            element["content"] = "[DIAGRAM/TABLE DETECTED]"
                            temp_elements.append(element)
                        elif box.label in ["Text", "Section-header", "List-item", "Title", "Caption"]:
                            crops_to_process.append(image.crop(box.bbox))
                            crop_indices.append(len(temp_elements))
                            temp_elements.append(element)

                    if crops_to_process:
                        BATCH_SIZE = 8
                        all_text_lines = []
                        for k in range(0, len(crops_to_process), BATCH_SIZE):
                            batch = crops_to_process[k : k + BATCH_SIZE]
                            all_text_lines.extend(rec_predictor(batch, det_predictor=det_predictor))
                        
                        total_chars = 0
                        for j, res in enumerate(all_text_lines):
                            txt = " ".join([l.text for l in res.text_lines]).strip()
                            temp_elements[crop_indices[j]]["content"] = txt
                            total_chars += len(txt)
                        
                        if total_chars < 50: raise Exception("Low Text")
                        page_data["elements"] = temp_elements
                    else:
                        raise Exception("No text blocks")

                except Exception:
                    # Fallback
                    ocr_result = rec_predictor([image], det_predictor=det_predictor)[0]
                    full_text = " ".join([l.text for l in ocr_result.text_lines])
                    page_data["elements"] = [{"type": "Full-Page-Text", "content": full_text}]

                response_data["pages"].append(page_data)
                del image
                gc.collect()

        elif engine == "tesseract":
             for i, img in enumerate(pil_images):
                text = pytesseract.image_to_string(img)
                response_data["pages"].append({"page": i+1, "elements": [{"type": "Raw-Text", "content": text}]})

    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)

    return JSONResponse(content=response_data)